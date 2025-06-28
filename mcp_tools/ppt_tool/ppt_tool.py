#!/usr/bin/env python3
"""
PowerPoint Presentation MCP Server
사용자가 PowerPoint 프레젠테이션을 생성하고 편집할 수 있는 도구를 제공합니다.
슬라이드 추가, 텍스트 편집, 이미지 삽입, 도형 추가, 서식 지정 등 다양한 기능을 포함합니다.
"""

import logging
import os
import sys
import time
import json
import base64
import io
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Any, Union, Tuple

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_THEME_COLOR
    from PIL import Image
except ImportError:
    print("필요한 라이브러리가 설치되어 있지 않습니다. 설치하려면: pip install python-pptx pillow")
    sys.exit(1)

from mcp.server.fastmcp import FastMCP

# --- 디버깅 로깅 설정 ---
# 이 스크립트가 별도 프로세스로 실행될 때의 오류를 추적하기 위함
# 프로젝트 루트에 ppt_tool_mcp.log 파일 생성
log_file_path = Path(__file__).resolve().parents[2] / "ppt_tool_mcp.log"
if os.path.exists(log_file_path):
    os.remove(log_file_path)

# 환경 변수로 로그 레벨 제어 (기본값: WARNING)
log_level = os.getenv("PPT_TOOL_LOG_LEVEL", "WARNING").upper()
log_level_int = getattr(logging, log_level, logging.WARNING)

logging.basicConfig(
    level=log_level_int,
    format="%(asctime)s [%(levelname)s] %(name)s: %(message)s",
    handlers=[logging.FileHandler(log_file_path),
              logging.StreamHandler(sys.stderr)],
)
logger = logging.getLogger(__name__)

# INFO 레벨 로그는 환경 변수가 DEBUG나 INFO로 설정된 경우에만 출력
if log_level_int <= logging.INFO:
    logger.info("PowerPoint Tool MCP 서버 프로세스 시작 (PID: %s)", os.getpid())
    logger.info("Python Executable: %s", sys.executable)
    logger.info("sys.path: %s", sys.path)
# --- 로깅 설정 끝 ---

# Create MCP Server
app = FastMCP(
    title="PowerPoint Presentation Server",
    description="A server for creating and editing PowerPoint presentations",
    version="1.0.0",
)

TRANSPORT = "stdio"

# 기본 설정
DEFAULT_SLIDE_WIDTH = 10  # 인치
DEFAULT_SLIDE_HEIGHT = 7.5  # 인치
DEFAULT_FONT_NAME = "Calibri"
DEFAULT_FONT_SIZE = 18  # 포인트
DEFAULT_TITLE_FONT_SIZE = 32  # 포인트
DEFAULT_TEMP_DIR = os.path.join(os.path.expanduser("~"), "ppt_tool_temp")

# 슬라이드 레이아웃 상수
LAYOUT_TITLE_SLIDE = 0  # 제목 슬라이드
LAYOUT_TITLE_CONTENT = 1  # 제목 및 내용
LAYOUT_SECTION_HEADER = 2  # 섹션 헤더
LAYOUT_TWO_CONTENT = 3  # 두 개의 내용
LAYOUT_COMPARISON = 4  # 비교
LAYOUT_TITLE_ONLY = 5  # 제목만
LAYOUT_BLANK = 6  # 빈 슬라이드
LAYOUT_CONTENT_CAPTION = 7  # 내용 및 캡션
LAYOUT_PICTURE_CAPTION = 8  # 그림 및 캡션

# 색상 상수
COLOR_BLACK = RGBColor(0, 0, 0)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_RED = RGBColor(255, 0, 0)
COLOR_GREEN = RGBColor(0, 255, 0)
COLOR_BLUE = RGBColor(0, 0, 255)
COLOR_YELLOW = RGBColor(255, 255, 0)
COLOR_CYAN = RGBColor(0, 255, 255)
COLOR_MAGENTA = RGBColor(255, 0, 255)

# 도형 상수
SHAPE_RECTANGLE = MSO_SHAPE.RECTANGLE
SHAPE_OVAL = MSO_SHAPE.OVAL
SHAPE_ROUNDED_RECTANGLE = MSO_SHAPE.ROUNDED_RECTANGLE
SHAPE_TRIANGLE = MSO_SHAPE.TRIANGLE
SHAPE_RIGHT_TRIANGLE = MSO_SHAPE.RIGHT_TRIANGLE
SHAPE_DIAMOND = MSO_SHAPE.DIAMOND
SHAPE_PENTAGON = MSO_SHAPE.PENTAGON
SHAPE_HEXAGON = MSO_SHAPE.HEXAGON
SHAPE_STAR = MSO_SHAPE.STAR_5_POINTS
SHAPE_ARROW = MSO_SHAPE.ARROW
SHAPE_LINE = MSO_SHAPE.LINE

# 텍스트 정렬 상수
ALIGN_LEFT = PP_ALIGN.LEFT
ALIGN_CENTER = PP_ALIGN.CENTER
ALIGN_RIGHT = PP_ALIGN.RIGHT
ALIGN_JUSTIFY = PP_ALIGN.JUSTIFY

# 수직 정렬 상수
VERTICAL_ALIGN_TOP = MSO_VERTICAL_ANCHOR.TOP
VERTICAL_ALIGN_MIDDLE = MSO_VERTICAL_ANCHOR.MIDDLE
VERTICAL_ALIGN_BOTTOM = MSO_VERTICAL_ANCHOR.BOTTOM


@dataclass
class SlideInfo:
    """슬라이드 정보를 담는 데이터 클래스"""
    index: int
    title: str = ""
    layout_type: int = LAYOUT_TITLE_CONTENT
    shapes_count: int = 0
    has_title: bool = False
    has_content: bool = False
    has_images: bool = False
    has_charts: bool = False
    has_tables: bool = False
    has_shapes: bool = False
    notes: str = ""


@dataclass
class PresentationInfo:
    """프레젠테이션 정보를 담는 데이터 클래스"""
    filename: str
    slides_count: int = 0
    slides: List[SlideInfo] = field(default_factory=list)
    author: str = ""
    title: str = ""
    subject: str = ""
    keywords: str = ""
    last_modified: str = ""
    created: str = ""
    slide_width: float = DEFAULT_SLIDE_WIDTH
    slide_height: float = DEFAULT_SLIDE_HEIGHT
    file_size: int = 0
    error: str = ""


@dataclass
class ShapeInfo:
    """도형 정보를 담는 데이터 클래스"""
    id: int
    name: str
    type: str
    left: float = 0.0
    top: float = 0.0
    width: float = 0.0
    height: float = 0.0
    text: str = ""
    rotation: float = 0.0
    z_order: int = 0
    has_text_frame: bool = False
    has_table: bool = False
    has_chart: bool = False
    is_placeholder: bool = False
    placeholder_type: str = ""


@dataclass
class TextFrameInfo:
    """텍스트 프레임 정보를 담는 데이터 클래스"""
    text: str
    font_name: str = DEFAULT_FONT_NAME
    font_size: int = DEFAULT_FONT_SIZE
    bold: bool = False
    italic: bool = False
    underline: bool = False
    color: str = "000000"  # RGB 색상 (16진수)
    alignment: str = "left"  # left, center, right, justify
    vertical_alignment: str = "top"  # top, middle, bottom
    margin_left: float = 0.1
    margin_right: float = 0.1
    margin_top: float = 0.05
    margin_bottom: float = 0.05


@dataclass
class TableInfo:
    """표 정보를 담는 데이터 클래스"""
    rows: int
    cols: int
    data: List[List[str]] = field(default_factory=list)
    header_row: bool = True
    style: str = "Medium Style 2 - Accent 1"


@dataclass
class ImageInfo:
    """이미지 정보를 담는 데이터 클래스"""
    filename: str = ""
    base64_data: str = ""
    left: float = 0.0
    top: float = 0.0
    width: float = 0.0
    height: float = 0.0
    crop_left: float = 0.0
    crop_right: float = 0.0
    crop_top: float = 0.0
    crop_bottom: float = 0.0
    brightness: float = 0.0
    contrast: float = 0.0
    transparency: float = 0.0


class PowerPointService:
    """PowerPoint 서비스 클래스"""

    def __init__(self):
        """PowerPoint 서비스 초기화"""
        # 임시 디렉토리 생성
        os.makedirs(DEFAULT_TEMP_DIR, exist_ok=True)

        # 현재 작업 중인 프레젠테이션 캐시
        self.presentations = {}  # 키: 파일 경로, 값: Presentation 객체

        logger.info("PowerPoint 서비스 초기화 완료")

    def _get_presentation(self, filename: str) -> Presentation:
        """
        파일 이름으로 프레젠테이션 객체를 가져옵니다.
        캐시에 없으면 파일에서 로드하거나 새로 생성합니다.

        Args:
            filename: 프레젠테이션 파일 경로

        Returns:
            Presentation: 프레젠테이션 객체
        """
        # 절대 경로로 변환
        filename = os.path.abspath(filename)

        # 캐시에 있으면 반환
        if filename in self.presentations:
            return self.presentations[filename]

        # 파일이 존재하면 로드, 아니면 새로 생성
        if os.path.exists(filename):
            try:
                prs = Presentation(filename)
                self.presentations[filename] = prs
                return prs
            except Exception as e:
                logger.error(f"프레젠테이션 로드 중 오류 발생: {e}")
                raise
        else:
            # 새 프레젠테이션 생성
            prs = Presentation()
            self.presentations[filename] = prs
            return prs

    def create_presentation(self, filename: str, title: str = "", author: str = "") -> PresentationInfo:
        """
        새 프레젠테이션을 생성합니다.

        Args:
            filename: 저장할 파일 경로
            title: 프레젠테이션 제목
            author: 작성자

        Returns:
            PresentationInfo: 프레젠테이션 정보
        """
        try:
            # 절대 경로로 변환
            filename = os.path.abspath(filename)

            # 디렉토리가 없으면 생성
            os.makedirs(os.path.dirname(filename), exist_ok=True)

            # 새 프레젠테이션 생성
            prs = Presentation()

            # 프레젠테이션 속성 설정
            if title:
                prs.core_properties.title = title
            if author:
                prs.core_properties.author = author
                prs.core_properties.last_modified_by = author

            # 현재 시간 설정
            now = datetime.now().strftime("%Y-%m-%dT%H:%M:%SZ")
            prs.core_properties.created = now
            prs.core_properties.modified = now

            # 캐시에 저장
            self.presentations[filename] = prs

            # 저장
            prs.save(filename)

            # 정보 반환
            return self.get_presentation_info(filename)

        except Exception as e:
            logger.error(f"프레젠테이션 생성 중 오류 발생: {e}")
            return PresentationInfo(
                filename=filename,
                error=f"프레젠테이션 생성 중 오류 발생: {str(e)}"
            )

    def open_presentation(self, filename: str) -> PresentationInfo:
        """
        기존 프레젠테이션을 엽니다.

        Args:
            filename: 프레젠테이션 파일 경로

        Returns:
            PresentationInfo: 프레젠테이션 정보
        """
        try:
            # 절대 경로로 변환
            filename = os.path.abspath(filename)

            # 파일이 존재하는지 확인
            if not os.path.exists(filename):
                return PresentationInfo(
                    filename=filename,
                    error=f"파일이 존재하지 않습니다: {filename}"
                )

            # 프레젠테이션 로드
            self._get_presentation(filename)

            # 정보 반환
            return self.get_presentation_info(filename)

        except Exception as e:
            logger.error(f"프레젠테이션 열기 중 오류 발생: {e}")
            return PresentationInfo(
                filename=filename,
                error=f"프레젠테이션 열기 중 오류 발생: {str(e)}"
            )

    def save_presentation(self, filename: str, new_filename: str = None) -> Dict[str, Any]:
        """
        프레젠테이션을 저장합니다.

        Args:
            filename: 현재 파일 경로
            new_filename: 새 파일 경로 (다른 이름으로 저장)

        Returns:
            Dict[str, Any]: 저장 결과
        """
        try:
            # 절대 경로로 변환
            filename = os.path.abspath(filename)

            # 프레젠테이션 가져오기
            prs = self._get_presentation(filename)

            # 다른 이름으로 저장
            if new_filename:
                new_filename = os.path.abspath(new_filename)
                # 디렉토리가 없으면 생성
                os.makedirs(os.path.dirname(new_filename), exist_ok=True)
                prs.save(new_filename)
                # 캐시 업데이트
                self.presentations[new_filename] = prs
                save_path = new_filename
            else:
                # 원래 파일에 저장
                prs.save(filename)
                save_path = filename

            return {
                "success": True,
                "filename": save_path,
                "message": f"프레젠테이션이 성공적으로 저장되었습니다: {save_path}"
            }

        except Exception as e:
            logger.error(f"프레젠테이션 저장 중 오류 발생: {e}")
            return {
                "success": False,
                "filename": filename,
                "error": f"프레젠테이션 저장 중 오류 발생: {str(e)}"
            }

    def close_presentation(self, filename: str) -> Dict[str, Any]:
        """
        프레젠테이션을 닫습니다.

        Args:
            filename: 프레젠테이션 파일 경로

        Returns:
            Dict[str, Any]: 닫기 결과
        """
        try:
            # 절대 경로로 변환
            filename = os.path.abspath(filename)

            # 캐시에서 제거
            if filename in self.presentations:
                del self.presentations[filename]
                return {
                    "success": True,
                    "message": f"프레젠테이션이 성공적으로 닫혔습니다: {filename}"
                }
            else:
                return {
                    "success": False,
                    "error": f"열려있는 프레젠테이션이 아닙니다: {filename}"
                }

        except Exception as e:
            logger.error(f"프레젠테이션 닫기 중 오류 발생: {e}")
            return {
                "success": False,
                "error": f"프레젠테이션 닫기 중 오류 발생: {str(e)}"
            }

    def get_presentation_info(self, filename: str) -> PresentationInfo:
        """
        프레젠테이션 정보를 가져옵니다.

        Args:
            filename: 프레젠테이션 파일 경로

        Returns:
            PresentationInfo: 프레젠테이션 정보
        """
        try:
            # 절대 경로로 변환
            filename = os.path.abspath(filename)

            # 프레젠테이션 가져오기
            prs = self._get_presentation(filename)

            # 기본 정보
            info = PresentationInfo(
                filename=filename,
                slides_count=len(prs.slides),
                author=prs.core_properties.author or "",
                title=prs.core_properties.title or "",
                subject=prs.core_properties.subject or "",
                keywords=prs.core_properties.keywords or "",
                last_modified=prs.core_properties.modified or "",
                created=prs.core_properties.created or "",
                slide_width=prs.slide_width.inches,
                slide_height=prs.slide_height.inches
            )

            # 파일 크기
            if os.path.exists(filename):
                info.file_size = os.path.getsize(filename)

            # 슬라이드 정보
            for i, slide in enumerate(prs.slides):
                slide_info = SlideInfo(
                    index=i,
                    layout_type=slide.slide_layout.index,
                    shapes_count=len(slide.shapes)
                )

                # 제목 확인
                if slide.shapes.title:
                    slide_info.has_title = True
                    slide_info.title = slide.shapes.title.text

                # 슬라이드 내용 확인
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        slide_info.has_content = True
                    if shape.shape_type == 13:  # 그림
                        slide_info.has_images = True
                    if shape.shape_type == 3:  # 차트
                        slide_info.has_charts = True
                    if shape.has_table:
                        slide_info.has_tables = True
                    if shape.shape_type in [1, 2, 5]:  # 자동 도형
                        slide_info.has_shapes = True

                # 노트 확인
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    slide_info.notes = slide.notes_slide.notes_text_frame.text

                info.slides.append(slide_info)

            return info

        except Exception as e:
            logger.error(f"프레젠테이션 정보 가져오기 중 오류 발생: {e}")
            return PresentationInfo(
                filename=filename,
                error=f"프레젠테이션 정보 가져오기 중 오류 발생: {str(e)}"
            )

    def add_slide(self, filename: str, layout_type: int = LAYOUT_TITLE_CONTENT, title: str = "", content: str = "") -> Dict[str, Any]:
        """
        슬라이드를 추가합니다.

        Args:
            filename: 프레젠테이션 파일 경로
            layout_type: 슬라이드 레이아웃 유형
            title: 슬라이드 제목
            content: 슬라이드 내용

        Returns:
            Dict[str, Any]: 추가 결과
        """
        try:
            # 프레젠테이션 가져오기
            prs = self._get_presentation(filename)

            # 레이아웃 유형 확인
            if layout_type < 0 or layout_type >= len(prs.slide_layouts):
                layout_type = LAYOUT_TITLE_CONTENT

            # 슬라이드 추가
            slide_layout = prs.slide_layouts[layout_type]
            slide = prs.slides.add_slide(slide_layout)

            # 제목 설정
            if slide.shapes.title and title:
                slide.shapes.title.text = title

            # 내용 설정
            if content:
                for shape in slide.placeholders:
                    if shape.placeholder_format.type == 1:  # 내용 placeholder
                        shape.text = content
                        break

            # 슬라이드 인덱스
            slide_index = len(prs.slides) - 1

            return {
                "success": True,
                "slide_index": slide_index,
                "message": f"슬라이드가 성공적으로 추가되었습니다 (인덱스: {slide_index})"
            }

        except Exception as e:
            logger.error(f"슬라이드 추가 중 오류 발생: {e}")
            return {
                "success": False,
                "error": f"슬라이드 추가 중 오류 발생: {str(e)}"
            }

    def delete_slide(self, filename: str, slide_index: int) -> Dict[str, Any]:
        """
        슬라이드를 삭제합니다.

        Args:
            filename: 프레젠테이션 파일 경로
            slide_index: 삭제할 슬라이드 인덱스

        Returns:
            Dict[str, Any]: 삭제 결과
        """
        try:
            # 프레젠테이션 가져오기
            prs = self._get_presentation(filename)

            # 슬라이드 인덱스 확인
            if slide_index < 0 or slide_index >= len(prs.slides):
                return {
                    "success": False,
                    "error": f"유효하지 않은 슬라이드 인덱스: {slide_index}"
                }

            # XML 구조에서 슬라이드 삭제
            slides = list(prs.slides._sldIdLst)
            prs.part.drop_rel(slides[slide_index].rId)
            del slides[slide_index]

            return {
                "success": True,
                "message": f"슬라이드가 성공적으로 삭제되었습니다 (인덱스: {slide_index})"
            }

        except Exception as e:
            logger.error(f"슬라이드 삭제 중 오류 발생: {e}")
            return {
                "success": False,
                "error": f"슬라이드 삭제 중 오류 발생: {str(e)}"
            }

    def update_slide_title(self, filename: str, slide_index: int, title: str) -> Dict[str, Any]:
        """
        슬라이드 제목을 업데이트합니다.

        Args:
            filename: 프레젠테이션 파일 경로
            slide_index: 슬라이드 인덱스
            title: 새 제목

        Returns:
            Dict[str, Any]: 업데이트 결과
        """
        try:
            # 프레젠테이션 가져오기
            prs = self._get_presentation(filename)

            # 슬라이드 인덱스 확인
            if slide_index < 0 or slide_index >= len(prs.slides):
                return {
                    "success": False,
                    "error": f"유효하지 않은 슬라이드 인덱스: {slide_index}"
                }

            # 슬라이드 가져오기
            slide = prs.slides[slide_index]

            # 제목 업데이트
            if slide.shapes.title:
                slide.shapes.title.text = title
                return {
                    "success": True,
                    "message": f"슬라이드 제목이 성공적으로 업데이트되었습니다 (인덱스: {slide_index})"
                }
            else:
                return {
                    "success": False,
                    "error": f"슬라이드에 제목 placeholder가 없습니다 (인덱스: {slide_index})"
                }

        except Exception as e:
            logger.error(f"슬라이드 제목 업데이트 중 오류 발생: {e}")
            return {
                "success": False,
                "error": f"슬라이드 제목 업데이트 중 오류 발생: {str(e)}"
            }

    def add_text_box(self, filename: str, slide_index: int, text: str, left: float = 1.0, top: float = 1.0, 
                     width: float = 4.0, height: float = 1.0, font_name: str = DEFAULT_FONT_NAME, 
                     font_size: int = DEFAULT_FONT_SIZE, bold: bool = False, italic: bool = False, 
                     color: str = "000000", alignment: str = "left") -> Dict[str, Any]:
        """
        텍스트 상자를 추가합니다.

        Args:
            filename: 프레젠테이션 파일 경로
            slide_index: 슬라이드 인덱스
            text: 텍스트 내용
            left: 왼쪽 위치 (인치)
            top: 위쪽 위치 (인치)
            width: 너비 (인치)
            height: 높이 (인치)
            font_name: 글꼴 이름
            font_size: 글꼴 크기 (포인트)
            bold: 굵게 여부
            italic: 기울임꼴 여부
            color: RGB 색상 (16진수)
            alignment: 정렬 (left, center, right, justify)

        Returns:
            Dict[str, Any]: 추가 결과
        """
        try:
            # 프레젠테이션 가져오기
            prs = self._get_presentation(filename)

            # 슬라이드 인덱스 확인
            if slide_index < 0 or slide_index >= len(prs.slides):
                return {
                    "success": False,
                    "error": f"유효하지 않은 슬라이드 인덱스: {slide_index}"
                }

            # 슬라이드 가져오기
            slide = prs.slides[slide_index]

            # 텍스트 상자 추가
            textbox = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(height)
            )

            # 텍스트 설정
            text_frame = textbox.text_frame
            text_frame.text = text

            # 단락 설정
            p = text_frame.paragraphs[0]

            # 정렬 설정
            if alignment.lower() == "center":
                p.alignment = ALIGN_CENTER
            elif alignment.lower() == "right":
                p.alignment = ALIGN_RIGHT
            elif alignment.lower() == "justify":
                p.alignment = ALIGN_JUSTIFY
            else:
                p.alignment = ALIGN_LEFT

            # 글꼴 설정
            run = p.runs[0]
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.italic = italic

            # 색상 설정
            if color:
                try:
                    r = int(color[0:2], 16)
                    g = int(color[2:4], 16)
                    b = int(color[4:6], 16)
                    run.font.color.rgb = RGBColor(r, g, b)
                except:
                    run.font.color.rgb = COLOR_BLACK

            return {
                "success": True,
                "shape_id": textbox.shape_id,
                "message": f"텍스트 상자가 성공적으로 추가되었습니다 (슬라이드: {slide_index}, ID: {textbox.shape_id})"
            }

        except Exception as e:
            logger.error(f"텍스트 상자 추가 중 오류 발생: {e}")
            return {
                "success": False,
                "error": f"텍스트 상자 추가 중 오류 발생: {str(e)}"
            }

    def add_shape(self, filename: str, slide_index: int, shape_type: int = SHAPE_RECTANGLE, 
                  left: float = 1.0, top: float = 1.0, width: float = 2.0, height: float = 1.0, 
                  text: str = "", fill_color: str = "", line_color: str = "000000", 
                  line_width: float = 1.0) -> Dict[str, Any]:
        """
        도형을 추가합니다.

        Args:
            filename: 프레젠테이션 파일 경로
            slide_index: 슬라이드 인덱스
            shape_type: 도형 유형
            left: 왼쪽 위치 (인치)
            top: 위쪽 위치 (인치)
            width: 너비 (인치)
            height: 높이 (인치)
            text: 도형 내 텍스트
            fill_color: 채우기 색상 (16진수)
            line_color: 선 색상 (16진수)
            line_width: 선 두께 (포인트)

        Returns:
            Dict[str, Any]: 추가 결과
        """
        try:
            # 프레젠테이션 가져오기
            prs = self._get_presentation(filename)

            # 슬라이드 인덱스 확인
            if slide_index < 0 or slide_index >= len(prs.slides):
                return {
                    "success": False,
                    "error": f"유효하지 않은 슬라이드 인덱스: {slide_index}"
                }

            # 슬라이드 가져오기
            slide = prs.slides[slide_index]

            # 도형 추가
            shape = slide.shapes.add_shape(
                shape_type, Inches(left), Inches(top), Inches(width), Inches(height)
            )

            # 채우기 색상 설정
            if fill_color:
                try:
                    r = int(fill_color[0:2], 16)
                    g = int(fill_color[2:4], 16)
                    b = int(fill_color[4:6], 16)
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(r, g, b)
                except:
                    pass

            # 선 색상 설정
            if line_color:
                try:
                    r = int(line_color[0:2], 16)
                    g = int(line_color[2:4], 16)
                    b = int(line_color[4:6], 16)
                    shape.line.color.rgb = RGBColor(r, g, b)
                except:
                    pass

            # 선 두께 설정
            shape.line.width = Pt(line_width)

            # 텍스트 설정
            if text:
                shape.text = text
                # 텍스트 중앙 정렬
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.alignment = ALIGN_CENTER
                shape.text_frame.vertical_anchor = VERTICAL_ALIGN_MIDDLE

            return {
                "success": True,
                "shape_id": shape.shape_id,
                "message": f"도형이 성공적으로 추가되었습니다 (슬라이드: {slide_index}, ID: {shape.shape_id})"
            }

        except Exception as e:
            logger.error(f"도형 추가 중 오류 발생: {e}")
            return {
                "success": False,
                "error": f"도형 추가 중 오류 발생: {str(e)}"
            }

    def add_image(self, filename: str, slide_index: int, image_path: str = None, base64_data: str = None, 
                  left: float = 1.0, top: float = 1.0, width: float = 4.0, height: float = 3.0) -> Dict[str, Any]:
        """
        이미지를 추가합니다.

        Args:
            filename: 프레젠테이션 파일 경로
            slide_index: 슬라이드 인덱스
            image_path: 이미지 파일 경로
            base64_data: Base64 인코딩된 이미지 데이터
            left: 왼쪽 위치 (인치)
            top: 위쪽 위치 (인치)
            width: 너비 (인치)
            height: 높이 (인치)

        Returns:
            Dict[str, Any]: 추가 결과
        """
        try:
            # 프레젠테이션 가져오기
            prs = self._get_presentation(filename)

            # 슬라이드 인덱스 확인
            if slide_index < 0 or slide_index >= len(prs.slides):
                return {
                    "success": False,
                    "error": f"유효하지 않은 슬라이드 인덱스: {slide_index}"
                }

            # 슬라이드 가져오기
            slide = prs.slides[slide_index]

            # 이미지 데이터 확인
            if not image_path and not base64_data:
                return {
                    "success": False,
                    "error": "이미지 파일 경로 또는 Base64 데이터가 필요합니다"
                }

            # 이미지 추가
            if image_path:
                # 파일 경로로 이미지 추가
                image_path = os.path.abspath(image_path)
                if not os.path.exists(image_path):
                    return {
                        "success": False,
                        "error": f"이미지 파일이 존재하지 않습니다: {image_path}"
                    }

                picture = slide.shapes.add_picture(
                    image_path, Inches(left), Inches(top), Inches(width), Inches(height)
                )
            else:
                # Base64 데이터로 이미지 추가
                try:
                    # Base64 디코딩
                    image_data = base64.b64decode(base64_data)

                    # 임시 파일로 저장
                    temp_image_path = os.path.join(DEFAULT_TEMP_DIR, f"temp_image_{int(time.time())}.png")
                    with open(temp_image_path, "wb") as f:
                        f.write(image_data)

                    # 이미지 추가
                    picture = slide.shapes.add_picture(
                        temp_image_path, Inches(left), Inches(top), Inches(width), Inches(height)
                    )

                    # 임시 파일 삭제
                    try:
                        os.remove(temp_image_path)
                    except:
                        pass

                except Exception as e:
                    return {
                        "success": False,
                        "error": f"Base64 이미지 데이터 처리 중 오류 발생: {str(e)}"
                    }

            return {
                "success": True,
                "shape_id": picture.shape_id,
                "message": f"이미지가 성공적으로 추가되었습니다 (슬라이드: {slide_index}, ID: {picture.shape_id})"
            }

        except Exception as e:
            logger.error(f"이미지 추가 중 오류 발생: {e}")
            return {
                "success": False,
                "error": f"이미지 추가 중 오류 발생: {str(e)}"
            }

    def add_table(self, filename: str, slide_index: int, rows: int, cols: int, 
                  data: List[List[str]] = None, left: float = 1.0, top: float = 1.0, 
                  width: float = 8.0, height: float = 4.0, style: str = None) -> Dict[str, Any]:
        """
        표를 추가합니다.

        Args:
            filename: 프레젠테이션 파일 경로
            slide_index: 슬라이드 인덱스
            rows: 행 수
            cols: 열 수
            data: 표 데이터 (2차원 리스트)
            left: 왼쪽 위치 (인치)
            top: 위쪽 위치 (인치)
            width: 너비 (인치)
            height: 높이 (인치)
            style: 표 스타일

        Returns:
            Dict[str, Any]: 추가 결과
        """
        try:
            # 프레젠테이션 가져오기
            prs = self._get_presentation(filename)

            # 슬라이드 인덱스 확인
            if slide_index < 0 or slide_index >= len(prs.slides):
                return {
                    "success": False,
                    "error": f"유효하지 않은 슬라이드 인덱스: {slide_index}"
                }

            # 슬라이드 가져오기
            slide = prs.slides[slide_index]

            # 행과 열 수 확인
            if rows < 1 or cols < 1:
                return {
                    "success": False,
                    "error": "행과 열 수는 1 이상이어야 합니다"
                }

            # 표 추가
            table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table

            # 데이터 설정
            if data:
                for i in range(min(rows, len(data))):
                    for j in range(min(cols, len(data[i]))):
                        cell = table.cell(i, j)
                        cell.text = str(data[i][j])

            # 스타일 설정
            if style and hasattr(table, "style"):
                try:
                    table.style = style
                except:
                    pass

            return {
                "success": True,
                "shape_id": table.parent.shape_id,
                "message": f"표가 성공적으로 추가되었습니다 (슬라이드: {slide_index}, 행: {rows}, 열: {cols})"
            }

        except Exception as e:
            logger.error(f"표 추가 중 오류 발생: {e}")
            return {
                "success": False,
                "error": f"표 추가 중 오류 발생: {str(e)}"
            }

    def delete_shape(self, filename: str, slide_index: int, shape_id: int) -> Dict[str, Any]:
        """
        도형을 삭제합니다.

        Args:
            filename: 프레젠테이션 파일 경로
            slide_index: 슬라이드 인덱스
            shape_id: 도형 ID

        Returns:
            Dict[str, Any]: 삭제 결과
        """
        try:
            # 프레젠테이션 가져오기
            prs = self._get_presentation(filename)

            # 슬라이드 인덱스 확인
            if slide_index < 0 or slide_index >= len(prs.slides):
                return {
                    "success": False,
                    "error": f"유효하지 않은 슬라이드 인덱스: {slide_index}"
                }

            # 슬라이드 가져오기
            slide = prs.slides[slide_index]

            # 도형 찾기 및 삭제
            shape_found = False
            for i, shape in enumerate(slide.shapes):
                if shape.shape_id == shape_id:
                    # XML 구조에서 도형 삭제
                    slide.shapes._spTree.remove(shape._element)
                    shape_found = True
                    break

            if shape_found:
                return {
                    "success": True,
                    "message": f"도형이 성공적으로 삭제되었습니다 (슬라이드: {slide_index}, ID: {shape_id})"
                }
            else:
                return {
                    "success": False,
                    "error": f"도형을 찾을 수 없습니다 (슬라이드: {slide_index}, ID: {shape_id})"
                }

        except Exception as e:
            logger.error(f"도형 삭제 중 오류 발생: {e}")
            return {
                "success": False,
                "error": f"도형 삭제 중 오류 발생: {str(e)}"
            }

    def extract_images(self, filename: str, output_folder: str = None) -> Dict[str, Any]:
        """
        프레젠테이션의 모든 이미지를 추출합니다.

        Args:
            filename: 프레젠테이션 파일 경로
            output_folder: 이미지를 저장할 폴더 경로 (기본값: 프레젠테이션 파일명_images)

        Returns:
            Dict[str, Any]: 추출 결과
        """
        try:
            # 절대 경로로 변환
            filename = os.path.abspath(filename)

            # 프레젠테이션 가져오기
            prs = self._get_presentation(filename)

            # 출력 폴더 설정
            if not output_folder:
                base_name = os.path.splitext(os.path.basename(filename))[0]
                output_folder = os.path.join(os.path.dirname(filename), f"{base_name}_images")

            # 폴더 생성
            os.makedirs(output_folder, exist_ok=True)

            # 이미지 추출
            extracted_images = []
            image_count = 0

            # 각 슬라이드에서 이미지 추출
            for slide_idx, slide in enumerate(prs.slides):
                for shape_idx, shape in enumerate(slide.shapes):
                    if shape.shape_type == 13:  # 그림
                        try:
                            # 이미지 파일 이름 생성
                            image_filename = f"slide_{slide_idx+1}_image_{shape_idx+1}.png"
                            image_path = os.path.join(output_folder, image_filename)

                            # 이미지 데이터 추출 및 저장
                            if hasattr(shape, 'image'):
                                # 이미지 데이터 추출
                                image_bytes = shape.image.blob

                                # 이미지 저장
                                with open(image_path, 'wb') as f:
                                    f.write(image_bytes)

                                # 이미지 정보 저장
                                image_info = ImageInfo(
                                    filename=image_filename,
                                    path=image_path,
                                    slide_index=slide_idx,
                                    shape_id=shape.shape_id,
                                    width=shape.width.inches,
                                    height=shape.height.inches,
                                    size_bytes=len(image_bytes)
                                )

                                extracted_images.append(image_info.__dict__)
                                image_count += 1
                        except Exception as e:
                            logger.warning(f"슬라이드 {slide_idx+1}의 이미지 추출 중 오류 발생: {e}")

            return {
                "success": True,
                "output_folder": output_folder,
                "images": extracted_images,
                "count": image_count,
                "message": f"{image_count}개의 이미지가 {output_folder}에 추출되었습니다."
            }

        except Exception as e:
            logger.error(f"이미지 추출 중 오류 발생: {e}")
            return {
                "success": False,
                "error": f"이미지 추출 중 오류 발생: {str(e)}"
            }

    def check_drm(self, filename: str) -> Dict[str, Any]:
        """
        프레젠테이션 파일의 DRM 보호 여부를 확인합니다.

        Args:
            filename: 프레젠테이션 파일 경로

        Returns:
            Dict[str, Any]: 확인 결과
        """
        try:
            # 절대 경로로 변환
            filename = os.path.abspath(filename)

            # 파일이 존재하는지 확인
            if not os.path.exists(filename):
                return {
                    "success": False,
                    "error": f"파일이 존재하지 않습니다: {filename}"
                }

            # 파일 확장자 확인
            _, ext = os.path.splitext(filename)
            ext = ext.lower()

            # DRM 확인 결과
            is_drm_protected = False
            drm_type = "none"
            can_open = True
            message = "DRM 보호가 적용되지 않은 파일입니다."

            # PPTX 파일 확인
            if ext == '.pptx':
                try:
                    # 파일 열기 시도
                    prs = Presentation(filename)

                    # 성공적으로 열렸으면 DRM이 없거나 python-pptx가 지원하는 형식
                    # 추가 검사: 슬라이드 접근 시도
                    slide_count = len(prs.slides)
                    if slide_count > 0:
                        # 첫 번째 슬라이드 접근 시도
                        _ = prs.slides[0]
                except Exception as e:
                    error_str = str(e).lower()
                    if "corrupt" in error_str or "invalid" in error_str or "permission" in error_str:
                        is_drm_protected = True
                        drm_type = "encryption"
                        can_open = False
                        message = f"파일이 암호화되어 있거나 손상되었습니다: {str(e)}"

            # PPT 파일 확인 (python-pptx는 .ppt 파일을 직접 지원하지 않음)
            elif ext == '.ppt':
                is_drm_protected = False  # 기본값
                drm_type = "unknown"
                can_open = False
                message = "PPT 형식은 python-pptx에서 직접 지원하지 않습니다. PPTX로 변환 후 사용하세요."

            return {
                "success": True,
                "filename": filename,
                "is_drm_protected": is_drm_protected,
                "drm_type": drm_type,
                "can_open": can_open,
                "message": message
            }

        except Exception as e:
            logger.error(f"DRM 확인 중 오류 발생: {e}")
            return {
                "success": False,
                "error": f"DRM 확인 중 오류 발생: {str(e)}"
            }

    def get_slide_shapes(self, filename: str, slide_index: int) -> Dict[str, Any]:
        """
        슬라이드의 모든 도형 정보를 가져옵니다.

        Args:
            filename: 프레젠테이션 파일 경로
            slide_index: 슬라이드 인덱스

        Returns:
            Dict[str, Any]: 도형 정보
        """
        try:
            # 프레젠테이션 가져오기
            prs = self._get_presentation(filename)

            # 슬라이드 인덱스 확인
            if slide_index < 0 or slide_index >= len(prs.slides):
                return {
                    "success": False,
                    "error": f"유효하지 않은 슬라이드 인덱스: {slide_index}"
                }

            # 슬라이드 가져오기
            slide = prs.slides[slide_index]

            # 도형 정보 수집
            shapes_info = []
            for shape in slide.shapes:
                shape_type = "unknown"
                if shape.shape_type == 1:
                    shape_type = "auto_shape"
                elif shape.shape_type == 2:
                    shape_type = "freeform"
                elif shape.shape_type == 3:
                    shape_type = "chart"
                elif shape.shape_type == 13:
                    shape_type = "picture"
                elif shape.shape_type == 14:
                    shape_type = "placeholder"
                elif shape.shape_type == 16:
                    shape_type = "group"
                elif shape.shape_type == 17:
                    shape_type = "table"
                elif shape.shape_type == 19:
                    shape_type = "text_box"

                shape_info = ShapeInfo(
                    id=shape.shape_id,
                    name=shape.name,
                    type=shape_type,
                    left=shape.left.inches,
                    top=shape.top.inches,
                    width=shape.width.inches,
                    height=shape.height.inches,
                    has_text_frame=shape.has_text_frame,
                    has_table=shape.has_table,
                    has_chart=hasattr(shape, "chart"),
                    is_placeholder=hasattr(shape, "is_placeholder") and shape.is_placeholder
                )

                # 텍스트 정보
                if shape.has_text_frame:
                    shape_info.text = shape.text

                # placeholder 유형
                if shape_info.is_placeholder and hasattr(shape, "placeholder_format"):
                    placeholder_types = {
                        0: "title",
                        1: "body",
                        2: "centered_title",
                        3: "subtitle",
                        4: "notes",
                        5: "header",
                        6: "footer",
                        7: "slide_number",
                        8: "date",
                        9: "slide_image",
                        10: "table",
                        11: "chart",
                        12: "diagram",
                        13: "media"
                    }
                    placeholder_type = shape.placeholder_format.type
                    shape_info.placeholder_type = placeholder_types.get(placeholder_type, str(placeholder_type))

                shapes_info.append(shape_info.__dict__)

            return {
                "success": True,
                "shapes": shapes_info,
                "count": len(shapes_info)
            }

        except Exception as e:
            logger.error(f"도형 정보 가져오기 중 오류 발생: {e}")
            return {
                "success": False,
                "error": f"도형 정보 가져오기 중 오류 발생: {str(e)}"
            }

    def update_shape_text(self, filename: str, slide_index: int, shape_id: int, text: str) -> Dict[str, Any]:
        """
        도형의 텍스트를 업데이트합니다.

        Args:
            filename: 프레젠테이션 파일 경로
            slide_index: 슬라이드 인덱스
            shape_id: 도형 ID
            text: 새 텍스트

        Returns:
            Dict[str, Any]: 업데이트 결과
        """
        try:
            # 프레젠테이션 가져오기
            prs = self._get_presentation(filename)

            # 슬라이드 인덱스 확인
            if slide_index < 0 or slide_index >= len(prs.slides):
                return {
                    "success": False,
                    "error": f"유효하지 않은 슬라이드 인덱스: {slide_index}"
                }

            # 슬라이드 가져오기
            slide = prs.slides[slide_index]

            # 도형 찾기 및 텍스트 업데이트
            shape_found = False
            for shape in slide.shapes:
                if shape.shape_id == shape_id:
                    if shape.has_text_frame:
                        shape.text_frame.text = text
                        shape_found = True
                        break
                    else:
                        return {
                            "success": False,
                            "error": f"도형에 텍스트 프레임이 없습니다 (슬라이드: {slide_index}, ID: {shape_id})"
                        }

            if shape_found:
                return {
                    "success": True,
                    "message": f"도형 텍스트가 성공적으로 업데이트되었습니다 (슬라이드: {slide_index}, ID: {shape_id})"
                }
            else:
                return {
                    "success": False,
                    "error": f"도형을 찾을 수 없습니다 (슬라이드: {slide_index}, ID: {shape_id})"
                }

        except Exception as e:
            logger.error(f"도형 텍스트 업데이트 중 오류 발생: {e}")
            return {
                "success": False,
                "error": f"도형 텍스트 업데이트 중 오류 발생: {str(e)}"
            }

    def update_shape_position(self, filename: str, slide_index: int, shape_id: int, 
                              left: float = None, top: float = None, 
                              width: float = None, height: float = None) -> Dict[str, Any]:
        """
        도형의 위치와 크기를 업데이트합니다.

        Args:
            filename: 프레젠테이션 파일 경로
            slide_index: 슬라이드 인덱스
            shape_id: 도형 ID
            left: 새 왼쪽 위치 (인치)
            top: 새 위쪽 위치 (인치)
            width: 새 너비 (인치)
            height: 새 높이 (인치)

        Returns:
            Dict[str, Any]: 업데이트 결과
        """
        try:
            # 프레젠테이션 가져오기
            prs = self._get_presentation(filename)

            # 슬라이드 인덱스 확인
            if slide_index < 0 or slide_index >= len(prs.slides):
                return {
                    "success": False,
                    "error": f"유효하지 않은 슬라이드 인덱스: {slide_index}"
                }

            # 슬라이드 가져오기
            slide = prs.slides[slide_index]

            # 도형 찾기 및 위치/크기 업데이트
            shape_found = False
            for shape in slide.shapes:
                if shape.shape_id == shape_id:
                    if left is not None:
                        shape.left = Inches(left)
                    if top is not None:
                        shape.top = Inches(top)
                    if width is not None:
                        shape.width = Inches(width)
                    if height is not None:
                        shape.height = Inches(height)
                    shape_found = True
                    break

            if shape_found:
                return {
                    "success": True,
                    "message": f"도형 위치/크기가 성공적으로 업데이트되었습니다 (슬라이드: {slide_index}, ID: {shape_id})"
                }
            else:
                return {
                    "success": False,
                    "error": f"도형을 찾을 수 없습니다 (슬라이드: {slide_index}, ID: {shape_id})"
                }

        except Exception as e:
            logger.error(f"도형 위치/크기 업데이트 중 오류 발생: {e}")
            return {
                "success": False,
                "error": f"도형 위치/크기 업데이트 중 오류 발생: {str(e)}"
            }

    def add_slide_notes(self, filename: str, slide_index: int, notes: str) -> Dict[str, Any]:
        """
        슬라이드 노트를 추가합니다.

        Args:
            filename: 프레젠테이션 파일 경로
            slide_index: 슬라이드 인덱스
            notes: 노트 내용

        Returns:
            Dict[str, Any]: 추가 결과
        """
        try:
            # 프레젠테이션 가져오기
            prs = self._get_presentation(filename)

            # 슬라이드 인덱스 확인
            if slide_index < 0 or slide_index >= len(prs.slides):
                return {
                    "success": False,
                    "error": f"유효하지 않은 슬라이드 인덱스: {slide_index}"
                }

            # 슬라이드 가져오기
            slide = prs.slides[slide_index]

            # 노트 추가
            if not slide.has_notes_slide:
                slide.notes_slide = prs.notes_master.notes_layouts[0].clone()

            slide.notes_slide.notes_text_frame.text = notes

            return {
                "success": True,
                "message": f"슬라이드 노트가 성공적으로 추가되었습니다 (슬라이드: {slide_index})"
            }

        except Exception as e:
            logger.error(f"슬라이드 노트 추가 중 오류 발생: {e}")
            return {
                "success": False,
                "error": f"슬라이드 노트 추가 중 오류 발생: {str(e)}"
            }

    def set_slide_background(self, filename: str, slide_index: int, color: str = None, 
                             image_path: str = None, base64_data: str = None) -> Dict[str, Any]:
        """
        슬라이드 배경을 설정합니다.

        Args:
            filename: 프레젠테이션 파일 경로
            slide_index: 슬라이드 인덱스
            color: 배경 색상 (16진수)
            image_path: 배경 이미지 파일 경로
            base64_data: Base64 인코딩된 배경 이미지 데이터

        Returns:
            Dict[str, Any]: 설정 결과
        """
        try:
            # 프레젠테이션 가져오기
            prs = self._get_presentation(filename)

            # 슬라이드 인덱스 확인
            if slide_index < 0 or slide_index >= len(prs.slides):
                return {
                    "success": False,
                    "error": f"유효하지 않은 슬라이드 인덱스: {slide_index}"
                }

            # 슬라이드 가져오기
            slide = prs.slides[slide_index]

            # 배경 설정
            if color:
                # 색상 배경 설정
                try:
                    r = int(color[0:2], 16)
                    g = int(color[2:4], 16)
                    b = int(color[4:6], 16)
                    background = slide.background
                    fill = background.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(r, g, b)
                    return {
                        "success": True,
                        "message": f"슬라이드 배경 색상이 성공적으로 설정되었습니다 (슬라이드: {slide_index})"
                    }
                except Exception as e:
                    return {
                        "success": False,
                        "error": f"배경 색상 설정 중 오류 발생: {str(e)}"
                    }
            elif image_path or base64_data:
                # 이미지 배경 설정 (직접적인 API가 없어 대안으로 전체 크기 이미지 추가)
                try:
                    # 슬라이드 크기
                    slide_width = prs.slide_width.inches
                    slide_height = prs.slide_height.inches

                    # 기존 배경 이미지 삭제
                    for shape in slide.shapes:
                        if hasattr(shape, "name") and "Background" in shape.name:
                            slide.shapes._spTree.remove(shape._element)

                    # 이미지 추가
                    if image_path:
                        # 파일 경로로 이미지 추가
                        image_path = os.path.abspath(image_path)
                        if not os.path.exists(image_path):
                            return {
                                "success": False,
                                "error": f"이미지 파일이 존재하지 않습니다: {image_path}"
                            }

                        picture = slide.shapes.add_picture(
                            image_path, 0, 0, Inches(slide_width), Inches(slide_height)
                        )
                    else:
                        # Base64 데이터로 이미지 추가
                        try:
                            # Base64 디코딩
                            image_data = base64.b64decode(base64_data)

                            # 임시 파일로 저장
                            temp_image_path = os.path.join(DEFAULT_TEMP_DIR, f"temp_bg_{int(time.time())}.png")
                            with open(temp_image_path, "wb") as f:
                                f.write(image_data)

                            # 이미지 추가
                            picture = slide.shapes.add_picture(
                                temp_image_path, 0, 0, Inches(slide_width), Inches(slide_height)
                            )

                            # 임시 파일 삭제
                            try:
                                os.remove(temp_image_path)
                            except:
                                pass

                        except Exception as e:
                            return {
                                "success": False,
                                "error": f"Base64 이미지 데이터 처리 중 오류 발생: {str(e)}"
                            }

                    # 이미지를 맨 뒤로 보내기
                    picture.name = "Background Image"
                    picture.z_order = 0

                    return {
                        "success": True,
                        "message": f"슬라이드 배경 이미지가 성공적으로 설정되었습니다 (슬라이드: {slide_index})"
                    }
                except Exception as e:
                    return {
                        "success": False,
                        "error": f"배경 이미지 설정 중 오류 발생: {str(e)}"
                    }
            else:
                return {
                    "success": False,
                    "error": "배경 색상 또는 이미지가 필요합니다"
                }

        except Exception as e:
            logger.error(f"슬라이드 배경 설정 중 오류 발생: {e}")
            return {
                "success": False,
                "error": f"슬라이드 배경 설정 중 오류 발생: {str(e)}"
            }

    def duplicate_slide(self, filename: str, slide_index: int) -> Dict[str, Any]:
        """
        슬라이드를 복제합니다.

        Args:
            filename: 프레젠테이션 파일 경로
            slide_index: 복제할 슬라이드 인덱스

        Returns:
            Dict[str, Any]: 복제 결과
        """
        try:
            # 프레젠테이션 가져오기
            prs = self._get_presentation(filename)

            # 슬라이드 인덱스 확인
            if slide_index < 0 or slide_index >= len(prs.slides):
                return {
                    "success": False,
                    "error": f"유효하지 않은 슬라이드 인덱스: {slide_index}"
                }

            # 슬라이드 가져오기
            source_slide = prs.slides[slide_index]

            # 슬라이드 레이아웃 가져오기
            slide_layout = source_slide.slide_layout

            # 새 슬라이드 추가
            new_slide = prs.slides.add_slide(slide_layout)

            # 도형 복제
            for shape in source_slide.shapes:
                if shape.shape_type == 14:  # placeholder
                    # placeholder 복제
                    placeholder = None
                    for placeholder_shape in new_slide.placeholders:
                        if placeholder_shape.placeholder_format.idx == shape.placeholder_format.idx:
                            placeholder = placeholder_shape
                            break

                    if placeholder and shape.has_text_frame and placeholder.has_text_frame:
                        placeholder.text = shape.text
                elif shape.shape_type == 13:  # picture
                    # 이미지는 복제할 수 없음 (API 제한)
                    pass
                elif shape.shape_type == 17:  # table
                    # 표는 복제할 수 없음 (API 제한)
                    pass
                elif shape.has_text_frame:
                    # 텍스트 상자 복제
                    new_shape = new_slide.shapes.add_textbox(
                        shape.left, shape.top, shape.width, shape.height
                    )
                    new_shape.text = shape.text

                    # 텍스트 서식 복제 (제한적)
                    if shape.text_frame.paragraphs and new_shape.text_frame.paragraphs:
                        for i, p in enumerate(shape.text_frame.paragraphs):
                            if i < len(new_shape.text_frame.paragraphs):
                                new_p = new_shape.text_frame.paragraphs[i]
                                new_p.alignment = p.alignment

            # 노트 복제
            if source_slide.has_notes_slide:
                if not new_slide.has_notes_slide:
                    new_slide.notes_slide = prs.notes_master.notes_layouts[0].clone()
                new_slide.notes_slide.notes_text_frame.text = source_slide.notes_slide.notes_text_frame.text

            # 새 슬라이드 인덱스
            new_slide_index = len(prs.slides) - 1

            return {
                "success": True,
                "new_slide_index": new_slide_index,
                "message": f"슬라이드가 성공적으로 복제되었습니다 (원본: {slide_index}, 새 슬라이드: {new_slide_index})"
            }

        except Exception as e:
            logger.error(f"슬라이드 복제 중 오류 발생: {e}")
            return {
                "success": False,
                "error": f"슬라이드 복제 중 오류 발생: {str(e)}"
            }

    def move_slide(self, filename: str, slide_index: int, new_position: int) -> Dict[str, Any]:
        """
        슬라이드 위치를 이동합니다.

        Args:
            filename: 프레젠테이션 파일 경로
            slide_index: 이동할 슬라이드 인덱스
            new_position: 새 위치 인덱스

        Returns:
            Dict[str, Any]: 이동 결과
        """
        try:
            # 프레젠테이션 가져오기
            prs = self._get_presentation(filename)

            # 슬라이드 수
            slides_count = len(prs.slides)

            # 인덱스 확인
            if slide_index < 0 or slide_index >= slides_count:
                return {
                    "success": False,
                    "error": f"유효하지 않은 슬라이드 인덱스: {slide_index}"
                }

            if new_position < 0 or new_position >= slides_count:
                return {
                    "success": False,
                    "error": f"유효하지 않은 새 위치 인덱스: {new_position}"
                }

            if slide_index == new_position:
                return {
                    "success": True,
                    "message": f"슬라이드가 이미 해당 위치에 있습니다 (인덱스: {slide_index})"
                }

            # XML 구조에서 슬라이드 이동
            slides = list(prs.slides._sldIdLst)
            slide_id = slides[slide_index]
            prs.slides._sldIdLst.remove(slide_id)
            prs.slides._sldIdLst.insert(new_position, slide_id)

            return {
                "success": True,
                "message": f"슬라이드가 성공적으로 이동되었습니다 (원래 위치: {slide_index}, 새 위치: {new_position})"
            }

        except Exception as e:
            logger.error(f"슬라이드 이동 중 오류 발생: {e}")
            return {
                "success": False,
                "error": f"슬라이드 이동 중 오류 발생: {str(e)}"
            }

    def set_presentation_properties(self, filename: str, title: str = None, author: str = None, 
                                    subject: str = None, keywords: str = None) -> Dict[str, Any]:
        """
        프레젠테이션 속성을 설정합니다.

        Args:
            filename: 프레젠테이션 파일 경로
            title: 제목
            author: 작성자
            subject: 주제
            keywords: 키워드

        Returns:
            Dict[str, Any]: 설정 결과
        """
        try:
            # 프레젠테이션 가져오기
            prs = self._get_presentation(filename)

            # 속성 설정
            if title is not None:
                prs.core_properties.title = title
            if author is not None:
                prs.core_properties.author = author
                prs.core_properties.last_modified_by = author
            if subject is not None:
                prs.core_properties.subject = subject
            if keywords is not None:
                prs.core_properties.keywords = keywords

            # 수정 시간 업데이트
            prs.core_properties.modified = datetime.now().strftime("%Y-%m-%dT%H:%M:%SZ")

            return {
                "success": True,
                "message": "프레젠테이션 속성이 성공적으로 설정되었습니다"
            }

        except Exception as e:
            logger.error(f"프레젠테이션 속성 설정 중 오류 발생: {e}")
            return {
                "success": False,
                "error": f"프레젠테이션 속성 설정 중 오류 발생: {str(e)}"
            }


# 싱글톤 인스턴스 생성
ppt_service = PowerPointService()


@app.tool()
def create_presentation(filename: str, title: str = "", author: str = "") -> dict:
    """
    새 프레젠테이션을 생성합니다.

    Args:
        filename: 저장할 파일 경로
        title: 프레젠테이션 제목
        author: 작성자

    Returns:
        dict: 생성 결과를 포함한 딕셔너리
    """
    try:
        result = ppt_service.create_presentation(filename, title, author)
        if result.error:
            return {"error": result.error}
        return {"result": result.__dict__}
    except Exception as e:
        logger.error(f"프레젠테이션 생성 중 오류 발생: {str(e)}")
        return {"error": f"프레젠테이션 생성 중 오류 발생: {str(e)}"}


@app.tool()
def open_presentation(filename: str) -> dict:
    """
    기존 프레젠테이션을 엽니다.

    Args:
        filename: 프레젠테이션 파일 경로

    Returns:
        dict: 열기 결과를 포함한 딕셔너리
    """
    try:
        result = ppt_service.open_presentation(filename)
        if result.error:
            return {"error": result.error}
        return {"result": result.__dict__}
    except Exception as e:
        logger.error(f"프레젠테이션 열기 중 오류 발생: {str(e)}")
        return {"error": f"프레젠테이션 열기 중 오류 발생: {str(e)}"}


@app.tool()
def save_presentation(filename: str, new_filename: str = None) -> dict:
    """
    프레젠테이션을 저장합니다.

    Args:
        filename: 현재 파일 경로
        new_filename: 새 파일 경로 (다른 이름으로 저장)

    Returns:
        dict: 저장 결과를 포함한 딕셔너리
    """
    try:
        result = ppt_service.save_presentation(filename, new_filename)
        if "error" in result:
            return {"error": result["error"]}
        return {"result": result}
    except Exception as e:
        logger.error(f"프레젠테이션 저장 중 오류 발생: {str(e)}")
        return {"error": f"프레젠테이션 저장 중 오류 발생: {str(e)}"}


@app.tool()
def close_presentation(filename: str) -> dict:
    """
    프레젠테이션을 닫습니다.

    Args:
        filename: 프레젠테이션 파일 경로

    Returns:
        dict: 닫기 결과를 포함한 딕셔너리
    """
    try:
        result = ppt_service.close_presentation(filename)
        if "error" in result:
            return {"error": result["error"]}
        return {"result": result}
    except Exception as e:
        logger.error(f"프레젠테이션 닫기 중 오류 발생: {str(e)}")
        return {"error": f"프레젠테이션 닫기 중 오류 발생: {str(e)}"}


@app.tool()
def get_presentation_info(filename: str) -> dict:
    """
    프레젠테이션 정보를 가져옵니다.

    Args:
        filename: 프레젠테이션 파일 경로

    Returns:
        dict: 프레젠테이션 정보를 포함한 딕셔너리
    """
    try:
        result = ppt_service.get_presentation_info(filename)
        if result.error:
            return {"error": result.error}
        return {"result": result.__dict__}
    except Exception as e:
        logger.error(f"프레젠테이션 정보 가져오기 중 오류 발생: {str(e)}")
        return {"error": f"프레젠테이션 정보 가져오기 중 오류 발생: {str(e)}"}


@app.tool()
def add_slide(filename: str, layout_type: int = LAYOUT_TITLE_CONTENT, title: str = "", content: str = "") -> dict:
    """
    슬라이드를 추가합니다.

    Args:
        filename: 프레젠테이션 파일 경로
        layout_type: 슬라이드 레이아웃 유형 (0: 제목 슬라이드, 1: 제목 및 내용, 5: 제목만, 6: 빈 슬라이드)
        title: 슬라이드 제목
        content: 슬라이드 내용

    Returns:
        dict: 추가 결과를 포함한 딕셔너리
    """
    try:
        result = ppt_service.add_slide(filename, layout_type, title, content)
        if "error" in result:
            return {"error": result["error"]}
        return {"result": result}
    except Exception as e:
        logger.error(f"슬라이드 추가 중 오류 발생: {str(e)}")
        return {"error": f"슬라이드 추가 중 오류 발생: {str(e)}"}


@app.tool()
def delete_slide(filename: str, slide_index: int) -> dict:
    """
    슬라이드를 삭제합니다.

    Args:
        filename: 프레젠테이션 파일 경로
        slide_index: 삭제할 슬라이드 인덱스

    Returns:
        dict: 삭제 결과를 포함한 딕셔너리
    """
    try:
        result = ppt_service.delete_slide(filename, slide_index)
        if "error" in result:
            return {"error": result["error"]}
        return {"result": result}
    except Exception as e:
        logger.error(f"슬라이드 삭제 중 오류 발생: {str(e)}")
        return {"error": f"슬라이드 삭제 중 오류 발생: {str(e)}"}


@app.tool()
def update_slide_title(filename: str, slide_index: int, title: str) -> dict:
    """
    슬라이드 제목을 업데이트합니다.

    Args:
        filename: 프레젠테이션 파일 경로
        slide_index: 슬라이드 인덱스
        title: 새 제목

    Returns:
        dict: 업데이트 결과를 포함한 딕셔너리
    """
    try:
        result = ppt_service.update_slide_title(filename, slide_index, title)
        if "error" in result:
            return {"error": result["error"]}
        return {"result": result}
    except Exception as e:
        logger.error(f"슬라이드 제목 업데이트 중 오류 발생: {str(e)}")
        return {"error": f"슬라이드 제목 업데이트 중 오류 발생: {str(e)}"}


@app.tool()
def add_text_box(filename: str, slide_index: int, text: str, left: float = 1.0, top: float = 1.0, 
                 width: float = 4.0, height: float = 1.0, font_name: str = DEFAULT_FONT_NAME, 
                 font_size: int = DEFAULT_FONT_SIZE, bold: bool = False, italic: bool = False, 
                 color: str = "000000", alignment: str = "left") -> dict:
    """
    텍스트 상자를 추가합니다.

    Args:
        filename: 프레젠테이션 파일 경로
        slide_index: 슬라이드 인덱스
        text: 텍스트 내용
        left: 왼쪽 위치 (인치)
        top: 위쪽 위치 (인치)
        width: 너비 (인치)
        height: 높이 (인치)
        font_name: 글꼴 이름
        font_size: 글꼴 크기 (포인트)
        bold: 굵게 여부
        italic: 기울임꼴 여부
        color: RGB 색상 (16진수)
        alignment: 정렬 (left, center, right, justify)

    Returns:
        dict: 추가 결과를 포함한 딕셔너리
    """
    try:
        result = ppt_service.add_text_box(filename, slide_index, text, left, top, width, height, 
                                          font_name, font_size, bold, italic, color, alignment)
        if "error" in result:
            return {"error": result["error"]}
        return {"result": result}
    except Exception as e:
        logger.error(f"텍스트 상자 추가 중 오류 발생: {str(e)}")
        return {"error": f"텍스트 상자 추가 중 오류 발생: {str(e)}"}


@app.tool()
def add_shape(filename: str, slide_index: int, shape_type: int = SHAPE_RECTANGLE, 
              left: float = 1.0, top: float = 1.0, width: float = 2.0, height: float = 1.0, 
              text: str = "", fill_color: str = "", line_color: str = "000000", 
              line_width: float = 1.0) -> dict:
    """
    도형을 추가합니다.

    Args:
        filename: 프레젠테이션 파일 경로
        slide_index: 슬라이드 인덱스
        shape_type: 도형 유형 (1: 직사각형, 2: 타원, 3: 둥근 직사각형, 4: 삼각형, 5: 마름모)
        left: 왼쪽 위치 (인치)
        top: 위쪽 위치 (인치)
        width: 너비 (인치)
        height: 높이 (인치)
        text: 도형 내 텍스트
        fill_color: 채우기 색상 (16진수)
        line_color: 선 색상 (16진수)
        line_width: 선 두께 (포인트)

    Returns:
        dict: 추가 결과를 포함한 딕셔너리
    """
    try:
        result = ppt_service.add_shape(filename, slide_index, shape_type, left, top, width, height, 
                                        text, fill_color, line_color, line_width)
        if "error" in result:
            return {"error": result["error"]}
        return {"result": result}
    except Exception as e:
        logger.error(f"도형 추가 중 오류 발생: {str(e)}")
        return {"error": f"도형 추가 중 오류 발생: {str(e)}"}


@app.tool()
def add_image(filename: str, slide_index: int, image_path: str = None, base64_data: str = None, 
              left: float = 1.0, top: float = 1.0, width: float = 4.0, height: float = 3.0) -> dict:
    """
    이미지를 추가합니다.

    Args:
        filename: 프레젠테이션 파일 경로
        slide_index: 슬라이드 인덱스
        image_path: 이미지 파일 경로
        base64_data: Base64 인코딩된 이미지 데이터
        left: 왼쪽 위치 (인치)
        top: 위쪽 위치 (인치)
        width: 너비 (인치)
        height: 높이 (인치)

    Returns:
        dict: 추가 결과를 포함한 딕셔너리
    """
    try:
        result = ppt_service.add_image(filename, slide_index, image_path, base64_data, left, top, width, height)
        if "error" in result:
            return {"error": result["error"]}
        return {"result": result}
    except Exception as e:
        logger.error(f"이미지 추가 중 오류 발생: {str(e)}")
        return {"error": f"이미지 추가 중 오류 발생: {str(e)}"}


@app.tool()
def add_table(filename: str, slide_index: int, rows: int, cols: int, 
              data: List[List[str]] = None, left: float = 1.0, top: float = 1.0, 
              width: float = 8.0, height: float = 4.0, style: str = None) -> dict:
    """
    표를 추가합니다.

    Args:
        filename: 프레젠테이션 파일 경로
        slide_index: 슬라이드 인덱스
        rows: 행 수
        cols: 열 수
        data: 표 데이터 (2차원 리스트)
        left: 왼쪽 위치 (인치)
        top: 위쪽 위치 (인치)
        width: 너비 (인치)
        height: 높이 (인치)
        style: 표 스타일

    Returns:
        dict: 추가 결과를 포함한 딕셔너리
    """
    try:
        result = ppt_service.add_table(filename, slide_index, rows, cols, data, left, top, width, height, style)
        if "error" in result:
            return {"error": result["error"]}
        return {"result": result}
    except Exception as e:
        logger.error(f"표 추가 중 오류 발생: {str(e)}")
        return {"error": f"표 추가 중 오류 발생: {str(e)}"}


@app.tool()
def delete_shape(filename: str, slide_index: int, shape_id: int) -> dict:
    """
    도형을 삭제합니다.

    Args:
        filename: 프레젠테이션 파일 경로
        slide_index: 슬라이드 인덱스
        shape_id: 도형 ID

    Returns:
        dict: 삭제 결과를 포함한 딕셔너리
    """
    try:
        result = ppt_service.delete_shape(filename, slide_index, shape_id)
        if "error" in result:
            return {"error": result["error"]}
        return {"result": result}
    except Exception as e:
        logger.error(f"도형 삭제 중 오류 발생: {str(e)}")
        return {"error": f"도형 삭제 중 오류 발생: {str(e)}"}


@app.tool()
def get_slide_shapes(filename: str, slide_index: int) -> dict:
    """
    슬라이드의 모든 도형 정보를 가져옵니다.

    Args:
        filename: 프레젠테이션 파일 경로
        slide_index: 슬라이드 인덱스

    Returns:
        dict: 도형 정보를 포함한 딕셔너리
    """
    try:
        result = ppt_service.get_slide_shapes(filename, slide_index)
        if "error" in result:
            return {"error": result["error"]}
        return {"result": result}
    except Exception as e:
        logger.error(f"도형 정보 가져오기 중 오류 발생: {str(e)}")
        return {"error": f"도형 정보 가져오기 중 오류 발생: {str(e)}"}


@app.tool()
def update_shape_text(filename: str, slide_index: int, shape_id: int, text: str) -> dict:
    """
    도형의 텍스트를 업데이트합니다.

    Args:
        filename: 프레젠테이션 파일 경로
        slide_index: 슬라이드 인덱스
        shape_id: 도형 ID
        text: 새 텍스트

    Returns:
        dict: 업데이트 결과를 포함한 딕셔너리
    """
    try:
        result = ppt_service.update_shape_text(filename, slide_index, shape_id, text)
        if "error" in result:
            return {"error": result["error"]}
        return {"result": result}
    except Exception as e:
        logger.error(f"도형 텍스트 업데이트 중 오류 발생: {str(e)}")
        return {"error": f"도형 텍스트 업데이트 중 오류 발생: {str(e)}"}


@app.tool()
def update_shape_position(filename: str, slide_index: int, shape_id: int, 
                          left: float = None, top: float = None, 
                          width: float = None, height: float = None) -> dict:
    """
    도형의 위치와 크기를 업데이트합니다.

    Args:
        filename: 프레젠테이션 파일 경로
        slide_index: 슬라이드 인덱스
        shape_id: 도형 ID
        left: 새 왼쪽 위치 (인치)
        top: 새 위쪽 위치 (인치)
        width: 새 너비 (인치)
        height: 새 높이 (인치)

    Returns:
        dict: 업데이트 결과를 포함한 딕셔너리
    """
    try:
        result = ppt_service.update_shape_position(filename, slide_index, shape_id, left, top, width, height)
        if "error" in result:
            return {"error": result["error"]}
        return {"result": result}
    except Exception as e:
        logger.error(f"도형 위치/크기 업데이트 중 오류 발생: {str(e)}")
        return {"error": f"도형 위치/크기 업데이트 중 오류 발생: {str(e)}"}


@app.tool()
def add_slide_notes(filename: str, slide_index: int, notes: str) -> dict:
    """
    슬라이드 노트를 추가합니다.

    Args:
        filename: 프레젠테이션 파일 경로
        slide_index: 슬라이드 인덱스
        notes: 노트 내용

    Returns:
        dict: 추가 결과를 포함한 딕셔너리
    """
    try:
        result = ppt_service.add_slide_notes(filename, slide_index, notes)
        if "error" in result:
            return {"error": result["error"]}
        return {"result": result}
    except Exception as e:
        logger.error(f"슬라이드 노트 추가 중 오류 발생: {str(e)}")
        return {"error": f"슬라이드 노트 추가 중 오류 발생: {str(e)}"}


@app.tool()
def set_slide_background(filename: str, slide_index: int, color: str = None, 
                         image_path: str = None, base64_data: str = None) -> dict:
    """
    슬라이드 배경을 설정합니다.

    Args:
        filename: 프레젠테이션 파일 경로
        slide_index: 슬라이드 인덱스
        color: 배경 색상 (16진수)
        image_path: 배경 이미지 파일 경로
        base64_data: Base64 인코딩된 배경 이미지 데이터

    Returns:
        dict: 설정 결과를 포함한 딕셔너리
    """
    try:
        result = ppt_service.set_slide_background(filename, slide_index, color, image_path, base64_data)
        if "error" in result:
            return {"error": result["error"]}
        return {"result": result}
    except Exception as e:
        logger.error(f"슬라이드 배경 설정 중 오류 발생: {str(e)}")
        return {"error": f"슬라이드 배경 설정 중 오류 발생: {str(e)}"}


@app.tool()
def duplicate_slide(filename: str, slide_index: int) -> dict:
    """
    슬라이드를 복제합니다.

    Args:
        filename: 프레젠테이션 파일 경로
        slide_index: 복제할 슬라이드 인덱스

    Returns:
        dict: 복제 결과를 포함한 딕셔너리
    """
    try:
        result = ppt_service.duplicate_slide(filename, slide_index)
        if "error" in result:
            return {"error": result["error"]}
        return {"result": result}
    except Exception as e:
        logger.error(f"슬라이드 복제 중 오류 발생: {str(e)}")
        return {"error": f"슬라이드 복제 중 오류 발생: {str(e)}"}


@app.tool()
def move_slide(filename: str, slide_index: int, new_position: int) -> dict:
    """
    슬라이드 위치를 이동합니다.

    Args:
        filename: 프레젠테이션 파일 경로
        slide_index: 이동할 슬라이드 인덱스
        new_position: 새 위치 인덱스

    Returns:
        dict: 이동 결과를 포함한 딕셔너리
    """
    try:
        result = ppt_service.move_slide(filename, slide_index, new_position)
        if "error" in result:
            return {"error": result["error"]}
        return {"result": result}
    except Exception as e:
        logger.error(f"슬라이드 이동 중 오류 발생: {str(e)}")
        return {"error": f"슬라이드 이동 중 오류 발생: {str(e)}"}


@app.tool()
def set_presentation_properties(filename: str, title: str = None, author: str = None, 
                                subject: str = None, keywords: str = None) -> dict:
    """
    프레젠테이션 속성을 설정합니다.

    Args:
        filename: 프레젠테이션 파일 경로
        title: 제목
        author: 작성자
        subject: 주제
        keywords: 키워드

    Returns:
        dict: 설정 결과를 포함한 딕셔너리
    """
    try:
        result = ppt_service.set_presentation_properties(filename, title, author, subject, keywords)
        if "error" in result:
            return {"error": result["error"]}
        return {"result": result}
    except Exception as e:
        logger.error(f"프레젠테이션 속성 설정 중 오류 발생: {str(e)}")
        return {"error": f"프레젠테이션 속성 설정 중 오류 발생: {str(e)}"}


@app.tool()
def extract_images(filename: str, output_folder: str = None) -> dict:
    """
    프레젠테이션의 모든 이미지를 추출합니다.

    Args:
        filename: 프레젠테이션 파일 경로
        output_folder: 이미지를 저장할 폴더 경로 (기본값: 프레젠테이션 파일명_images)

    Returns:
        dict: 추출 결과를 포함한 딕셔너리
    """
    try:
        result = ppt_service.extract_images(filename, output_folder)
        if "error" in result:
            return {"error": result["error"]}
        return {"result": result}
    except Exception as e:
        logger.error(f"이미지 추출 중 오류 발생: {str(e)}")
        return {"error": f"이미지 추출 중 오류 발생: {str(e)}"}


@app.tool()
def check_drm(filename: str) -> dict:
    """
    프레젠테이션 파일의 DRM 보호 여부를 확인합니다.

    Args:
        filename: 프레젠테이션 파일 경로

    Returns:
        dict: 확인 결과를 포함한 딕셔너리
    """
    try:
        result = ppt_service.check_drm(filename)
        if "error" in result:
            return {"error": result["error"]}
        return {"result": result}
    except Exception as e:
        logger.error(f"DRM 확인 중 오류 발생: {str(e)}")
        return {"error": f"DRM 확인 중 오류 발생: {str(e)}"}


@app.tool()
def get_tool_info() -> dict:
    """
    PowerPoint 도구 정보를 반환합니다.

    Returns:
        dict: 도구 정보를 포함한 딕셔너리
    """
    try:
        return {
            "result": {
                "name": "PowerPoint Presentation Tool",
                "description": "PowerPoint 프레젠테이션을 생성하고 편집하는 도구",
                "tools": [
                    {"name": "create_presentation", "description": "새 프레젠테이션을 생성합니다"},
                    {"name": "extract_images", "description": "프레젠테이션의 모든 이미지를 추출합니다"},
                    {"name": "check_drm", "description": "프레젠테이션 파일의 DRM 보호 여부를 확인합니다"},
                    {"name": "open_presentation", "description": "기존 프레젠테이션을 엽니다"},
                    {"name": "save_presentation", "description": "프레젠테이션을 저장합니다"},
                    {"name": "close_presentation", "description": "프레젠테이션을 닫습니다"},
                    {"name": "get_presentation_info", "description": "프레젠테이션 정보를 가져옵니다"},
                    {"name": "add_slide", "description": "슬라이드를 추가합니다"},
                    {"name": "delete_slide", "description": "슬라이드를 삭제합니다"},
                    {"name": "update_slide_title", "description": "슬라이드 제목을 업데이트합니다"},
                    {"name": "add_text_box", "description": "텍스트 상자를 추가합니다"},
                    {"name": "add_shape", "description": "도형을 추가합니다"},
                    {"name": "add_image", "description": "이미지를 추가합니다"},
                    {"name": "add_table", "description": "표를 추가합니다"},
                    {"name": "delete_shape", "description": "도형을 삭제합니다"},
                    {"name": "get_slide_shapes", "description": "슬라이드의 모든 도형 정보를 가져옵니다"},
                    {"name": "update_shape_text", "description": "도형의 텍스트를 업데이트합니다"},
                    {"name": "update_shape_position", "description": "도형의 위치와 크기를 업데이트합니다"},
                    {"name": "add_slide_notes", "description": "슬라이드 노트를 추가합니다"},
                    {"name": "set_slide_background", "description": "슬라이드 배경을 설정합니다"},
                    {"name": "duplicate_slide", "description": "슬라이드를 복제합니다"},
                    {"name": "move_slide", "description": "슬라이드 위치를 이동합니다"},
                    {"name": "set_presentation_properties", "description": "프레젠테이션 속성을 설정합니다"}
                ],
                "usage_examples": [
                    {"command": "create_presentation('presentation.pptx', '제목', '작성자')", "description": "새 프레젠테이션 생성"},
                    {"command": "add_slide('presentation.pptx', 1, '슬라이드 제목', '슬라이드 내용')", "description": "슬라이드 추가"},
                    {"command": "add_text_box('presentation.pptx', 0, '텍스트 내용', 1.0, 1.0, 4.0, 1.0)", "description": "텍스트 상자 추가"},
                    {"command": "add_image('presentation.pptx', 0, image_path='image.jpg')", "description": "이미지 추가"},
                    {"command": "save_presentation('presentation.pptx')", "description": "프레젠테이션 저장"}
                ],
                "dependencies": {
                    "required": [
                        "python-pptx - PowerPoint 파일 조작 라이브러리",
                        "pillow - 이미지 처리 라이브러리"
                    ],
                    "installation": "pip install python-pptx pillow"
                },
                "constants": {
                    "slide_layouts": {
                        "LAYOUT_TITLE_SLIDE": 0,
                        "LAYOUT_TITLE_CONTENT": 1,
                        "LAYOUT_SECTION_HEADER": 2,
                        "LAYOUT_TWO_CONTENT": 3,
                        "LAYOUT_COMPARISON": 4,
                        "LAYOUT_TITLE_ONLY": 5,
                        "LAYOUT_BLANK": 6,
                        "LAYOUT_CONTENT_CAPTION": 7,
                        "LAYOUT_PICTURE_CAPTION": 8
                    },
                    "shapes": {
                        "SHAPE_RECTANGLE": 1,
                        "SHAPE_OVAL": 2,
                        "SHAPE_ROUNDED_RECTANGLE": 3,
                        "SHAPE_TRIANGLE": 4,
                        "SHAPE_RIGHT_TRIANGLE": 5,
                        "SHAPE_DIAMOND": 6,
                        "SHAPE_PENTAGON": 7,
                        "SHAPE_HEXAGON": 8,
                        "SHAPE_STAR": 9,
                        "SHAPE_ARROW": 10,
                        "SHAPE_LINE": 11
                    }
                }
            }
        }

    except Exception as e:
        return {"error": f"도구 정보를 가져오는 중 오류 발생: {str(e)}"}


if __name__ == "__main__":
    try:
        logger.info("FastMCP app.run() 호출 시작...")
        app.run(transport=TRANSPORT)
        logger.info("FastMCP app.run() 정상 종료.")
    except Exception as e:
        logger.error("ppt_tool.py 스크립트 최상위에서 처리되지 않은 예외 발생", exc_info=True)
        raise
